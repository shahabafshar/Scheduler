"""
Create PowerPoint presentation from PRESENTATION.md using MY ECPE TEMPLATE.pptx
Enhanced version with proper layouts and detailed content for standalone deliverable.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# Paths
TEMPLATE_PATH = "MY ECPE TEMPLATE.pptx"
OUTPUT_PATH = "documents/final/Real_Time_Scheduling_Simulator_Presentation.pptx"
FIGURES_PATH = "documents/final/figures"
SCREENSHOTS_PATH = "user_guide/screenshots"


def add_text_with_bullets(text_frame, items, font_size=14, start_fresh=True):
    """Add bulleted items to a text frame"""
    if start_fresh:
        text_frame.clear()
    for i, item in enumerate(items):
        if i == 0 and start_fresh:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        # Handle nested items (tuples of (text, level))
        if isinstance(item, tuple):
            p.text = item[0]
            p.level = item[1]
        else:
            p.text = item
            p.level = 0
        p.font.size = Pt(font_size)


def add_image_safe(slide, img_path, left, top, width=None, height=None):
    """Add image if it exists, return True if added"""
    if os.path.exists(img_path):
        if width:
            slide.shapes.add_picture(img_path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(img_path, left, top, height=height)
        else:
            slide.shapes.add_picture(img_path, left, top)
        return True
    else:
        print(f"Warning: Image not found: {img_path}")
        return False


def add_image_to_placeholder(slide, placeholder, img_path):
    """Add image to a placeholder area, fitting within its bounds"""
    if not os.path.exists(img_path):
        print(f"Warning: Image not found: {img_path}")
        return False

    # Get placeholder position and size
    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height

    # Add picture at placeholder position
    slide.shapes.add_picture(img_path, left, top, width=width, height=height)

    # Hide the placeholder (optional - it may have default text)
    try:
        placeholder.text_frame.clear()
    except:
        pass

    return True


def create_presentation():
    # Load template
    prs = Presentation(TEMPLATE_PATH)

    # Clear existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # Get layouts
    title_layout = prs.slide_layouts[0]           # Title Slide
    content_layout = prs.slide_layouts[1]         # Title and Content
    two_content_layout = prs.slide_layouts[3]     # Two Content (LEFT + RIGHT)
    comparison_layout = prs.slide_layouts[4]      # Comparison
    title_only_layout = prs.slide_layouts[5]      # Title Only

    # ========================================================================
    # SLIDE 1: Title
    # ========================================================================
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "Real-Time Scheduling Simulator"
    slide.placeholders[1].text = "Server-Based Algorithms for Mixed Periodic-Aperiodic Workloads"

    # Presenter info text box
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(6), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Shahab Afshar"
    p.font.size = Pt(20)
    p.font.bold = True

    for text in ["CPR E 458/558: Real-Time Systems, Fall 2024",
                 "Department of Electrical and Computer Engineering",
                 "Iowa State University",
                 "Instructor: Dr. G. Manimaran"]:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(14)

    # ========================================================================
    # SLIDE 2: Problem Statement (Content only - no image)
    # ========================================================================
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Problem Statement"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # Main statement
    p = tf.paragraphs[0]
    p.text = "Real-time systems must satisfy timing constraints in addition to functional correctness."
    p.font.size = Pt(18)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = ""

    # Automotive example
    p = tf.add_paragraph()
    p.text = "Real-World Example: Automotive Engine Control Unit (ECU)"
    p.font.size = Pt(16)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "An ECU must handle both predictable and unpredictable workloads:"
    p.font.size = Pt(14)

    examples = [
        ("Periodic tasks: Engine control (10ms), ABS monitoring (5ms), fuel injection (20ms)", 1),
        ("Aperiodic tasks: Driver button presses, diagnostic requests, error handling", 1),
        ("Challenge: How to guarantee BOTH periodic AND aperiodic tasks meet deadlines?", 1),
    ]
    for text, level in examples:
        p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = ""

    # The Gap
    p = tf.add_paragraph()
    p.text = "The Gap in Current Tools"
    p.font.size = Pt(16)
    p.font.bold = True

    gaps = [
        ("Theoretical analysis provides utilization bounds but doesn't show temporal behavior", 1),
        ("Manual schedule construction is tedious and error-prone", 1),
        ("RTOS testing requires significant development effort before exploring alternatives", 1),
        ("Engineers need to explore algorithm behavior BEFORE committing to implementation", 1),
    ]
    for text, level in gaps:
        p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(14)

    # ========================================================================
    # SLIDE 3: Research Questions (Two Content: text + image)
    # ========================================================================
    slide = prs.slides.add_slide(two_content_layout)
    slide.shapes.title.text = "Research Questions"

    # Left content - Research Questions
    left_content = slide.placeholders[1]
    tf = left_content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "This project investigates three research questions:"
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = ""

    rqs = [
        ("RQ1: Faithful Implementation", "Can a discrete-event simulator faithfully implement server capacity management policies (Polling, Deferrable, Sporadic)?"),
        ("RQ2: Visualization Effectiveness", "Does interactive visualization (Gantt charts, metrics) help users understand algorithm differences?"),
        ("RQ3: Parameter Exploration", "Can parameter exploration (varying server capacity Cs and period Ps) reveal optimal configurations?"),
    ]

    for title, desc in rqs:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(12)

        p = tf.add_paragraph()
        p.text = ""

    # Right content - System diagram
    right_content = slide.placeholders[2]
    img_path = os.path.join(FIGURES_PATH, "system_diagram.png")
    add_image_to_placeholder(slide, right_content, img_path)

    # ========================================================================
    # SLIDE 4: Server Algorithms (Two Content: text + image)
    # ========================================================================
    slide = prs.slides.add_slide(two_content_layout)
    slide.shapes.title.text = "Solution: Four Server-Based Scheduling Algorithms"

    # Left content - Algorithm descriptions
    left_content = slide.placeholders[1]
    tf = left_content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Server-based scheduling dedicates CPU capacity to aperiodic tasks while protecting periodic task deadlines."
    p.font.size = Pt(12)

    p = tf.add_paragraph()
    p.text = ""

    servers = [
        ("Polling Server", "Capacity is LOST if no aperiodic tasks are ready when server executes. Simple but wasteful."),
        ("Deferrable Server", "Capacity is PRESERVED until end of server period. Better response but affects periodic schedulability."),
        ("Sporadic Server", "Capacity replenishment scheduled at t+Ps after consumption. Best response while maintaining RMS guarantees."),
        ("Background Scheduler", "Aperiodic tasks run ONLY during CPU idle time. Simplest but worst response time."),
    ]

    for name, desc in servers:
        p = tf.add_paragraph()
        p.text = name
        p.font.size = Pt(13)
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(11)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Sources: Sprunt et al. (1989), Lehoczky et al. (1987)"
    p.font.size = Pt(10)
    p.font.italic = True

    # Right content - Server comparison image
    right_content = slide.placeholders[2]
    img_path = os.path.join(FIGURES_PATH, "server_comparison.jpg")
    add_image_to_placeholder(slide, right_content, img_path)

    # ========================================================================
    # SLIDE 5: Platform Features (Two Content: text + screenshot)
    # ========================================================================
    slide = prs.slides.add_slide(two_content_layout)
    slide.shapes.title.text = "Interactive Simulation Platform"

    # Left content
    left_content = slide.placeholders[1]
    tf = left_content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "A web-based simulator built with Python and Streamlit providing:"
    p.font.size = Pt(13)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Core Features"
    p.font.bold = True
    p.font.size = Pt(14)

    features = [
        "Discrete-event simulation with exact algorithm implementation",
        "Interactive Gantt charts showing task execution and server capacity events",
        "Real-time parameter exploration via sliders (Cs, Ps, task parameters)",
        "21 preset configurations from course materials for quick experimentation",
        "Schedulability analysis with RMS, EDF, and DMS utilization bounds",
    ]
    for f in features:
        p = tf.add_paragraph()
        p.text = f
        p.level = 1
        p.font.size = Pt(11)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Supported Algorithms"
    p.font.bold = True
    p.font.size = Pt(14)

    algos = [
        "Basic: RMS, EDF, DMS, LLF (Liu & Layland, 1973)",
        "Server-Based: Polling, Deferrable, Sporadic, Background",
        "Advanced: Precedence-constrained, Overload handling, HVDF",
    ]
    for a in algos:
        p = tf.add_paragraph()
        p.text = a
        p.level = 1
        p.font.size = Pt(11)

    # Right content - Full layout screenshot
    right_content = slide.placeholders[2]
    img_path = os.path.join(SCREENSHOTS_PATH, "part1-getting-started/part1-02-full-layout.png")
    add_image_to_placeholder(slide, right_content, img_path)

    # ========================================================================
    # SLIDE 6: Architecture (Two Content: text + diagram)
    # ========================================================================
    slide = prs.slides.add_slide(two_content_layout)
    slide.shapes.title.text = "Software Architecture"

    # Left content
    left_content = slide.placeholders[1]
    tf = left_content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Layered Design with Clear Separation of Concerns"
    p.font.bold = True
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "The architecture separates the simulation engine from visualization and UI, enabling independent testing and extension."
    p.font.size = Pt(12)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Technology Stack"
    p.font.bold = True
    p.font.size = Pt(14)

    stack = [
        ("Core Logic (Python 3.10+)", "Pure simulation engine with no UI dependencies. Implements all scheduling algorithms and produces ScheduleResult objects."),
        ("Visualization (Plotly)", "Creates interactive Gantt charts, priority timelines, and metrics dashboards from ScheduleResult data."),
        ("Web UI (Streamlit)", "Provides the interactive dashboard with configuration panels, result displays, and export functionality."),
        ("Data Handling (Pandas)", "Manages task tables, CSV/JSON export, and data transformation for visualization."),
    ]

    for name, desc in stack:
        p = tf.add_paragraph()
        p.text = name
        p.font.size = Pt(12)
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(10)

    # Right content - Architecture diagram
    right_content = slide.placeholders[2]
    img_path = os.path.join(FIGURES_PATH, "layered_architecture.png")
    add_image_to_placeholder(slide, right_content, img_path)

    # ========================================================================
    # SLIDE 7: Template Method Pattern (Two Content: text + class diagram)
    # ========================================================================
    slide = prs.slides.add_slide(two_content_layout)
    slide.shapes.title.text = "Design Pattern: Template Method"

    # Left content
    left_content = slide.placeholders[1]
    tf = left_content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Key Insight"
    p.font.bold = True
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = "All schedulers share the same simulation loop (time advancement, ready queue management, preemption handling). Only priority assignment and task selection differ."
    p.font.size = Pt(12)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Implementation"
    p.font.bold = True
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = "SchedulerBase implements the complete simulate() method. Subclasses override only two methods:"
    p.font.size = Pt(12)

    methods = [
        "assign_priorities(): How to assign priorities to tasks",
        "get_next_task(ready_queue): Which task to run next",
    ]
    for m in methods:
        p = tf.add_paragraph()
        p.text = m
        p.level = 1
        p.font.size = Pt(11)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Server Capacity Management"
    p.font.bold = True
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = "The key difference between servers is what happens when the server has CPU priority but no aperiodic tasks are ready:"
    p.font.size = Pt(11)

    diffs = [
        "Polling: server_remaining = 0 (capacity lost)",
        "Deferrable: return False (capacity preserved until period end)",
        "Sporadic: Schedule replenishment at current_time + Ps",
    ]
    for d in diffs:
        p = tf.add_paragraph()
        p.text = d
        p.level = 1
        p.font.size = Pt(11)

    # Right content - Class hierarchy
    right_content = slide.placeholders[2]
    img_path = os.path.join(FIGURES_PATH, "class_hierarchy.png")
    add_image_to_placeholder(slide, right_content, img_path)

    # ========================================================================
    # SLIDE 8: RQ1 Results (Two Content: text + Gantt chart)
    # ========================================================================
    slide = prs.slides.add_slide(two_content_layout)
    slide.shapes.title.text = "Results: Correctness Validation (RQ1)"

    # Left content
    left_content = slide.placeholders[1]
    tf = left_content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "RQ1: Can the simulator faithfully implement server capacity management?"
    p.font.bold = True
    p.font.size = Pt(13)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Verification Method"
    p.font.bold = True
    p.font.size = Pt(13)

    p = tf.add_paragraph()
    p.text = "Each server type produces characteristic events in the timeline. Correctness is verified by observing expected event patterns:"
    p.font.size = Pt(11)

    verifications = [
        ("Polling Server", "capacity_lost events appear when server executes with no aperiodic tasks ready"),
        ("Deferrable Server", "deferred events show capacity preservation; no capacity_lost events"),
        ("Sporadic Server", "replenish events scheduled exactly at consumption_time + Ps"),
        ("Background Scheduler", "Aperiodic execution only during idle intervals (no periodic tasks ready)"),
    ]

    for name, desc in verifications:
        p = tf.add_paragraph()
        p.text = f"{name}: ✓"
        p.font.size = Pt(12)
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(10)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Answer: YES - The simulator correctly implements all four server algorithms as verified by behavioral event analysis."
    p.font.size = Pt(11)
    p.font.bold = True

    # Right content - Gantt chart
    right_content = slide.placeholders[2]
    img_path = os.path.join(SCREENSHOTS_PATH, "part3-server-algorithms/part3-polling-02-gantt-chart.png")
    add_image_to_placeholder(slide, right_content, img_path)

    # ========================================================================
    # SLIDE 9: RQ3 Results (Two Content: text + Gantt)
    # ========================================================================
    slide = prs.slides.add_slide(two_content_layout)
    slide.shapes.title.text = "Results: Parameter Sensitivity (RQ3)"

    # Left content
    left_content = slide.placeholders[1]
    tf = left_content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "RQ3: Can parameter exploration reveal optimal configurations?"
    p.font.bold = True
    p.font.size = Pt(13)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Experiment Setup"
    p.font.bold = True
    p.font.size = Pt(13)

    p = tf.add_paragraph()
    p.text = "Workload: Two periodic tasks P1(C=1, P=10), P2(C=1, P=15) and one aperiodic task A1(C=8, arrives at t=0). Server period Ps=5 fixed, varying capacity Cs."
    p.font.size = Pt(11)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Results: Server Capacity Effect on Response Time"
    p.font.bold = True
    p.font.size = Pt(13)

    results = [
        "Cs=2: Response Time = 17 time units (4 replenishment cycles needed)",
        "Cs=4: Response Time = 9 time units (2 replenishment cycles needed)",
        "Cs=8: Response Time = 8 time units (1 replenishment cycle needed)",
    ]
    for r in results:
        p = tf.add_paragraph()
        p.text = r
        p.level = 1
        p.font.size = Pt(11)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Key Finding"
    p.font.bold = True
    p.font.size = Pt(13)

    p = tf.add_paragraph()
    p.text = "Increasing server capacity (Cs) from 2 to 8 reduces aperiodic response time by 53% (17→8). This insight cannot be obtained from utilization formulas alone."
    p.font.size = Pt(11)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Answer: YES - Parameter exploration via sliders enables rapid discovery of optimal server configurations."
    p.font.size = Pt(11)
    p.font.bold = True

    # Right content - Sporadic Gantt
    right_content = slide.placeholders[2]
    img_path = os.path.join(SCREENSHOTS_PATH, "part3-server-algorithms/part3-sporadic-01-gantt.png")
    add_image_to_placeholder(slide, right_content, img_path)

    # ========================================================================
    # SLIDE 10: RQ2 Results (Two Content: text + Gantt features)
    # ========================================================================
    slide = prs.slides.add_slide(two_content_layout)
    slide.shapes.title.text = "Results: Visualization Effectiveness (RQ2)"

    # Left content
    left_content = slide.placeholders[1]
    tf = left_content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "RQ2: Does interactive visualization help users understand algorithm differences?"
    p.font.bold = True
    p.font.size = Pt(13)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Visualization Features Implemented"
    p.font.bold = True
    p.font.size = Pt(13)

    features = [
        ("Color-coded task bars", "Periodic tasks (blue), Aperiodic tasks (orange), Server capacity (green) are visually distinct"),
        ("Deadline markers", "Red triangles mark absolute deadlines for each task instance"),
        ("Arrival markers", "Green circles indicate task arrival times"),
        ("Server event labels", "Events like 'replenish', 'deferred', 'capacity_lost' labeled directly on timeline"),
        ("Hover tooltips", "Hovering over any bar shows task ID, remaining computation time, and event details"),
        ("Interactive zoom/pan", "Users can zoom into specific time regions for detailed analysis"),
    ]

    for name, desc in features:
        p = tf.add_paragraph()
        p.text = name
        p.font.size = Pt(12)
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(10)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Answer: YES - Users can visually compare when capacity is lost (Polling) vs. preserved (Deferrable) vs. replenished (Sporadic)."
    p.font.size = Pt(11)
    p.font.bold = True

    # Right content - Deferrable Gantt (different server to show variety)
    right_content = slide.placeholders[2]
    img_path = os.path.join(SCREENSHOTS_PATH, "part3-server-algorithms/part3-deferrable-01-gantt.png")
    add_image_to_placeholder(slide, right_content, img_path)

    # ========================================================================
    # SLIDE 11: Conclusions (Content only)
    # ========================================================================
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Conclusions"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Research Question Answers"
    p.font.bold = True
    p.font.size = Pt(16)

    answers = [
        ("RQ1: Faithful Implementation ✓", "The simulator correctly implements capacity management for all four server algorithms, verified by observing expected event patterns (capacity_lost, deferred, replenish)."),
        ("RQ2: Visualization Effectiveness ✓", "Interactive Gantt charts clearly show when capacity is lost vs. preserved vs. replenished. Users can visually compare algorithm behavior side-by-side."),
        ("RQ3: Parameter Exploration ✓", "Varying server capacity Cs from 2→8 demonstrated response time reduction from 17→8 time units, enabling discovery of optimal configurations."),
    ]

    for title, desc in answers:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(12)

        p = tf.add_paragraph()
        p.text = ""

    p = tf.add_paragraph()
    p.text = "Key Contributions"
    p.font.bold = True
    p.font.size = Pt(16)

    contributions = [
        "Open-source simulator implementing 4 server-based scheduling algorithms",
        "Interactive Gantt charts displaying capacity management events",
        "21 preset configurations from course literature for verified correctness",
        "Parameter exploration via sliders for optimal configuration discovery",
        "Schedulability analysis with RMS, EDF, and DMS utilization bounds",
    ]
    for i, c in enumerate(contributions):
        p = tf.add_paragraph()
        p.text = f"{i+1}. {c}"
        p.level = 1
        p.font.size = Pt(12)

    # ========================================================================
    # SLIDE 12: Limitations & Future Work (Content only)
    # ========================================================================
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Limitations & Future Work"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Current Limitations"
    p.font.bold = True
    p.font.size = Pt(16)

    limitations = [
        ("Single-processor only", "The simulator currently supports uniprocessor scheduling. Multi-core scenarios with task migration and partitioning are not yet implemented."),
        ("No resource contention modeling", "Priority Inheritance Protocol (PIP) and Priority Ceiling Protocol (PCP) are implemented in the codebase but not yet integrated into the main simulation loop."),
        ("Limited user validation", "RQ2 assessment is based on feature implementation. A formal user study would provide stronger evidence of visualization effectiveness."),
        ("Simplified execution model", "The simulator does not model I/O delays, cache effects, interrupt latency, or context switch overhead beyond basic counting."),
    ]

    for name, desc in limitations:
        p = tf.add_paragraph()
        p.text = name
        p.font.size = Pt(13)
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(11)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Future Directions"
    p.font.bold = True
    p.font.size = Pt(16)

    future = [
        "Multi-core extension with partitioned and global scheduling algorithms",
        "Resource protocol integration (PIP/PCP) for shared resource scenarios",
        "RTOS trace import to compare simulation with FreeRTOS/Zephyr execution",
        "Formal verification using model checking for schedulability guarantees",
    ]
    for f in future:
        p = tf.add_paragraph()
        p.text = f
        p.level = 1
        p.font.size = Pt(12)

    # ========================================================================
    # SLIDE 13: Learning Achieved (Content only)
    # ========================================================================
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Learning Achieved"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Technical Skills Developed"
    p.font.bold = True
    p.font.size = Pt(16)

    skills = [
        ("Real-Time Systems Concepts", "Deep understanding of server-based scheduling (Polling, Deferrable, Sporadic), RMS utilization bounds n(2^(1/n)-1), schedulability analysis, and capacity management trade-offs."),
        ("Software Engineering Patterns", "Applied Template Method pattern for extensible scheduler architecture, Strategy pattern for composable priority policies, and event-driven simulation design."),
        ("Tools & Technologies", "Python for discrete-event simulation, Streamlit for rapid web UI development, Plotly for interactive visualizations, Pandas for data handling."),
    ]

    for name, desc in skills:
        p = tf.add_paragraph()
        p.text = name
        p.font.size = Pt(14)
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(11)

        p = tf.add_paragraph()
        p.text = ""

    p = tf.add_paragraph()
    p.text = "Key Lessons Learned"
    p.font.bold = True
    p.font.size = Pt(16)

    lessons = [
        "Simulation reveals temporal behavior that utilization formulas cannot show",
        "Visualization is essential for understanding scheduling dynamics",
        "Parameter exploration enables informed design decisions before implementation",
        "Modular architecture (Template Method) dramatically simplifies adding new algorithms",
    ]
    for i, l in enumerate(lessons):
        p = tf.add_paragraph()
        p.text = f"{i+1}. {l}"
        p.level = 1
        p.font.size = Pt(12)

    # ========================================================================
    # SLIDE 14: References (Content only)
    # ========================================================================
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "References"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    refs = [
        '[1] B. Sprunt, L. Sha, and J. Lehoczky, "Aperiodic task scheduling for hard-real-time systems," Real-Time Systems, vol. 1, no. 1, pp. 27-60, 1989.',
        '[2] J. P. Lehoczky, L. Sha, and J. K. Strosnider, "Enhanced aperiodic responsiveness in hard real-time environments," in Proc. IEEE Real-Time Systems Symposium, pp. 261-270, 1987.',
        '[3] M. Spuri and G. Buttazzo, "Scheduling aperiodic tasks in dynamic priority systems," Real-Time Systems, vol. 10, no. 2, pp. 179-210, 1996.',
        '[4] C. L. Liu and J. W. Layland, "Scheduling algorithms for multiprogramming in a hard-real-time environment," Journal of the ACM, vol. 20, no. 1, pp. 46-61, 1973.',
    ]

    for i, ref in enumerate(refs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = ref
        p.font.size = Pt(12)
        p.space_after = Pt(14)

    # ========================================================================
    # SLIDE 15: Thank You
    # ========================================================================
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "Thank You"
    slide.placeholders[1].text = "Questions & Discussion"

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(6), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Shahab Afshar"
    p.font.size = Pt(18)
    p.font.bold = True

    for text in ["PhD Student, Computer Engineering",
                 "Department of Electrical and Computer Engineering",
                 "Iowa State University",
                 "",
                 "Project Repository: Available upon request"]:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(14)

    # Save
    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
